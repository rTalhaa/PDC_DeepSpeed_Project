from pathlib import Path
import json

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output" / "slides"
ASSETS = OUT / "assets"
OUT.mkdir(parents=True, exist_ok=True)
ASSETS.mkdir(parents=True, exist_ok=True)

DEEPSPEED_RESULTS = ROOT / "DeepSpeed" / "results" / "metrics_zero2_fair.json"
GPU_SUMMARY = ROOT / "DeepSpeed" / "results" / "gpu_summary_zero2_fair.csv"
TRAINER_STATE = ROOT / "DeepSpeed" / "outputs" / "zero2_fair_run" / "trainer_state.json"

NAVY = RGBColor(18, 32, 47)
INK = RGBColor(31, 39, 50)
MUTED = RGBColor(92, 104, 118)
TEAL = RGBColor(29, 146, 139)
GOLD = RGBColor(232, 167, 59)
RED = RGBColor(201, 85, 75)
BLUE = RGBColor(61, 102, 170)
GREEN = RGBColor(76, 154, 91)
LIGHT = RGBColor(244, 247, 250)
LINE = RGBColor(214, 222, 230)
WHITE = RGBColor(255, 255, 255)


qlora = {
    "hardware": "1 x Tesla T4",
    "method": "Single-GPU QLoRA adapter tuning",
    "trainable_params_m": 8.798,
    "total_params_m": 502.831,
    "trainable_pct": 1.7497,
    "epochs": 5,
    "steps": 565,
    "runtime": 2829.545,
    "samples_s": 1.590,
    "steps_s": 0.200,
    "train_loss": 0.448,
    "eval_loss": 0.823,
    "peak_allocated_gb": 1.66,
    "peak_reserved_gb": 3.56,
    "flops": 1.204e15,
}

with open(DEEPSPEED_RESULTS, "r", encoding="utf-8") as f:
    ds_metrics = json.load(f)
with open(TRAINER_STATE, "r", encoding="utf-8") as f:
    trainer_state = json.load(f)

train = ds_metrics["train_metrics"]
evalm = ds_metrics["eval_metrics"]
deepspeed = {
    "hardware": "2 x RTX A4000",
    "method": "Dual-GPU DeepSpeed ZeRO-2 full fine-tuning",
    "trainable_params_m": train["trainable_params"] / 1e6,
    "total_params_m": train["total_params"] / 1e6,
    "trainable_pct": train["trainable_percentage"],
    "epochs": train["epoch"],
    "steps": trainer_state["global_step"],
    "runtime": train["train_runtime"],
    "samples_s": train["train_samples_per_second"],
    "steps_s": train["train_steps_per_second"],
    "train_loss": train["train_loss"],
    "eval_loss": evalm["eval_loss"],
    "peak_allocated_gb": train["rank_peak_memory_allocated_gb"],
    "peak_reserved_gb": train["rank_peak_memory_reserved_gb"],
    "flops": train["total_flos"],
}


def save_bar_chart(filename, title, labels, values, ylabel, colors=None, annotate_fmt="{:.2f}"):
    fig, ax = plt.subplots(figsize=(7.2, 4.2), dpi=180)
    colors = colors or ["#1d928b", "#e8a73b"]
    bars = ax.bar(labels, values, color=colors, width=0.56)
    ax.set_title(title, loc="left", fontsize=15, fontweight="bold", color="#1f2732")
    ax.set_ylabel(ylabel, color="#5c6876")
    ax.grid(axis="y", color="#d6dee6", linewidth=0.8)
    ax.set_axisbelow(True)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color("#d6dee6")
    ax.spines["bottom"].set_color("#d6dee6")
    ax.tick_params(colors="#5c6876")
    ymax = max(values) * 1.18 if max(values) else 1
    ax.set_ylim(0, ymax)
    for b, v in zip(bars, values):
        ax.text(
            b.get_x() + b.get_width() / 2,
            b.get_height() + ymax * 0.025,
            annotate_fmt.format(v),
            ha="center",
            va="bottom",
            fontsize=10,
            color="#1f2732",
            fontweight="bold",
        )
    fig.tight_layout()
    path = ASSETS / filename
    fig.savefig(path, transparent=False, facecolor="white")
    plt.close(fig)
    return path


def save_grouped_loss():
    fig, ax = plt.subplots(figsize=(7.2, 4.2), dpi=180)
    x = np.arange(2)
    width = 0.34
    ax.bar(x - width / 2, [qlora["train_loss"], qlora["eval_loss"]], width, label="QLoRA", color="#1d928b")
    ax.bar(x + width / 2, [deepspeed["train_loss"], deepspeed["eval_loss"]], width, label="DeepSpeed ZeRO-2", color="#e8a73b")
    ax.set_title("Final Loss Comparison", loc="left", fontsize=15, fontweight="bold", color="#1f2732")
    ax.set_xticks(x, ["Training loss", "Evaluation loss"])
    ax.set_ylabel("Loss")
    ax.grid(axis="y", color="#d6dee6", linewidth=0.8)
    ax.set_axisbelow(True)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.legend(frameon=False)
    for container in ax.containers:
        ax.bar_label(container, fmt="%.3f", fontsize=9, padding=3)
    fig.tight_layout()
    path = ASSETS / "loss_comparison.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def save_training_curve():
    history = [h for h in trainer_state["log_history"] if "loss" in h]
    steps = [h["step"] for h in history]
    losses = [h["loss"] for h in history]
    fig, ax = plt.subplots(figsize=(7.2, 4.2), dpi=180)
    ax.plot(steps, losses, marker="o", linewidth=2.5, color="#3d66aa", label="Training loss")
    ax.axhline(deepspeed["eval_loss"], linestyle="--", color="#c9554b", label=f"Eval loss {deepspeed['eval_loss']:.3f}")
    ax.set_title("DeepSpeed ZeRO-2 Training Trend", loc="left", fontsize=15, fontweight="bold", color="#1f2732")
    ax.set_xlabel("Global step")
    ax.set_ylabel("Loss")
    ax.grid(color="#d6dee6", linewidth=0.8)
    ax.set_axisbelow(True)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.legend(frameon=False)
    fig.tight_layout()
    path = ASSETS / "deepspeed_loss_curve.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def save_gpu_chart():
    df = pd.read_csv(GPU_SUMMARY)
    labels = [f"GPU {int(i)}" for i in df["index"]]
    fig, axes = plt.subplots(1, 2, figsize=(9.2, 3.8), dpi=180)
    axes[0].bar(labels, df["avg_gpu_util_percent"], color="#1d928b")
    axes[0].set_title("Average GPU Utilization", loc="left", fontsize=13, fontweight="bold")
    axes[0].set_ylabel("Percent")
    axes[1].bar(labels, df["avg_memory_used_mib"] / 1024, color="#e8a73b")
    axes[1].set_title("Average Memory Used", loc="left", fontsize=13, fontweight="bold")
    axes[1].set_ylabel("GB")
    for ax in axes:
        ax.grid(axis="y", color="#d6dee6", linewidth=0.8)
        ax.set_axisbelow(True)
        for spine in ["top", "right"]:
            ax.spines[spine].set_visible(False)
    for ax in axes:
        for container in ax.containers:
            ax.bar_label(container, fmt="%.1f", fontsize=9, padding=3)
    fig.tight_layout()
    path = ASSETS / "gpu_telemetry.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


runtime_chart = save_bar_chart(
    "runtime_comparison.png",
    "Training Runtime",
    ["QLoRA", "DeepSpeed"],
    [qlora["runtime"], deepspeed["runtime"]],
    "Seconds",
    annotate_fmt="{:,.0f}",
)
throughput_chart = save_bar_chart(
    "throughput_comparison.png",
    "Training Throughput",
    ["QLoRA", "DeepSpeed"],
    [qlora["samples_s"], deepspeed["samples_s"]],
    "Samples / second",
)
params_chart = save_bar_chart(
    "trainable_params.png",
    "Trainable Parameter Share",
    ["QLoRA", "DeepSpeed"],
    [qlora["trainable_pct"], deepspeed["trainable_pct"]],
    "Percent of model updated",
    annotate_fmt="{:.1f}%",
)
memory_chart = save_bar_chart(
    "memory_comparison.png",
    "Peak Allocated GPU Memory",
    ["QLoRA", "DeepSpeed"],
    [qlora["peak_allocated_gb"], deepspeed["peak_allocated_gb"]],
    "GB",
)
loss_chart = save_grouped_loss()
curve_chart = save_training_curve()
gpu_chart = save_gpu_chart()


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.62), Inches(0.35), Inches(8.9), Inches(0.58))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.64), Inches(0.92), Inches(10.7), Inches(0.35))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Aptos"
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def add_footer(slide, idx):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(7.05), Inches(12.1), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(11.6), Inches(7.09), Inches(1.0), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"{idx:02d}"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def bullet_box(slide, x, y, w, h, bullets, font_size=17, color=INK, gap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10 if gap else 4)
    return box


def add_card(slide, x, y, w, h, title, body=None, fill=LIGHT, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.17), Inches(w - 0.42), Inches(0.36))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if body:
        b = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.62), Inches(w - 0.42), Inches(h - 0.75))
        bp = b.text_frame.paragraphs[0]
        bp.text = body
        bp.font.name = "Aptos"
        bp.font.size = Pt(12)
        bp.font.color.rgb = MUTED
        b.text_frame.word_wrap = True
    return shape


def add_table(slide, x, y, w, h, rows, col_widths=None, font_size=10.5):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (LIGHT if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE if r == 0 else INK
                p.font.bold = r == 0
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    return table_shape


def add_picture(slide, path, x, y, w, h):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def add_metric(slide, x, y, label, value, color):
    add_card(slide, x, y, 2.8, 1.15, label, "", fill=WHITE, accent=color)
    box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.55), Inches(2.3), Inches(0.38))
    p = box.text_frame.paragraphs[0]
    p.text = value
    p.font.name = "Aptos Display"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color


slides = []

# 1 Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
title = slide.shapes.add_textbox(Inches(0.75), Inches(0.78), Inches(11.8), Inches(1.8))
p = title.text_frame.paragraphs[0]
p.text = "Efficient Fine-Tuning of Code Language Models Under Hardware Constraints"
p.font.name = "Aptos Display"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = NAVY
subtitle = slide.shapes.add_textbox(Inches(0.78), Inches(2.65), Inches(10.9), Inches(0.8))
sp = subtitle.text_frame.paragraphs[0]
sp.text = "Comparative study of single-GPU QLoRA and dual-GPU DeepSpeed ZeRO-2"
sp.font.name = "Aptos"
sp.font.size = Pt(20)
sp.font.color.rgb = MUTED
add_metric(slide, 0.82, 4.25, "Model", "Qwen2.5-0.5B", TEAL)
add_metric(slide, 3.95, 4.25, "Dataset", "CodeAlpaca 20K", GOLD)
add_metric(slide, 7.08, 4.25, "Split", "900 / 100", BLUE)
add_metric(slide, 10.21, 4.25, "Focus", "Efficiency", GREEN)
footer = slide.shapes.add_textbox(Inches(0.78), Inches(6.55), Inches(8.5), Inches(0.35))
fp = footer.text_frame.paragraphs[0]
fp.text = "Talha Rashid, Bilal Fayyaz | Department of Computer Science"
fp.font.size = Pt(12)
fp.font.color.rgb = MUTED
add_footer(slide, 1)

# 2 Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Problem Introduction", "Fine-tuning code LLMs is useful, but hardware limits shape what is actually possible.")
bullet_box(
    slide,
    0.75,
    1.55,
    5.55,
    4.55,
    [
        "Full-parameter LLM fine-tuning is constrained by GPU memory, runtime, and distributed setup complexity.",
        "Student and research environments often rely on a single rented or notebook GPU.",
        "Parameter-efficient methods reduce memory pressure, while distributed methods improve throughput when multiple GPUs are available.",
        "The central question is how these two efficiency strategies trade off loss, speed, memory, and trainable parameter coverage.",
    ],
)
add_card(slide, 7.0, 1.55, 5.1, 1.25, "Single-GPU Constraint", "Google Colab Tesla T4 with QLoRA adapter tuning.", fill=LIGHT, accent=TEAL)
add_card(slide, 7.0, 3.12, 5.1, 1.25, "Multi-GPU Opportunity", "RunPod with two RTX A4000 GPUs using DeepSpeed ZeRO-2.", fill=LIGHT, accent=GOLD)
add_card(slide, 7.0, 4.69, 5.1, 1.25, "Comparison Target", "Same model family, same 1,000-sample CodeAlpaca subset, same 900/100 split.", fill=LIGHT, accent=BLUE)
add_footer(slide, 2)

# 3 Literature review
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Literature Review", "The base paper frames the study as a practical systems comparison rather than a new algorithm.")
add_card(slide, 0.75, 1.45, 3.85, 2.0, "Base Paper", "Efficient Fine-Tuning of Code Language Models Under Hardware Constraints compares QLoRA and DeepSpeed ZeRO-2 on Qwen2.5-0.5B-Instruct.", fill=WHITE, accent=BLUE)
add_card(slide, 4.9, 1.45, 3.55, 2.0, "LoRA and QLoRA", "LoRA trains low-rank adapter matrices. QLoRA adds 4-bit quantization to make adapter training feasible on limited memory.", fill=WHITE, accent=TEAL)
add_card(slide, 8.75, 1.45, 3.55, 2.0, "DeepSpeed ZeRO", "ZeRO partitions optimizer states and gradients across data-parallel workers, reducing memory redundancy in distributed training.", fill=WHITE, accent=GOLD)
add_card(slide, 0.75, 4.1, 3.85, 1.55, "Code Instruction Tuning", "CodeAlpaca 20K supplies prompt-completion examples for programming tasks.", fill=LIGHT, accent=GREEN)
add_card(slide, 4.9, 4.1, 3.55, 1.55, "Transformers and PEFT", "Hugging Face tooling supports model loading, tokenization, adapters, trainer metrics, and reproducible experiment tracking.", fill=LIGHT, accent=RED)
add_card(slide, 8.75, 4.1, 3.55, 1.55, "Research Gap", "A classroom-scale, metric-backed comparison helps decide when adapter tuning is enough and when distributed full fine-tuning is justified.", fill=LIGHT, accent=BLUE)
add_footer(slide, 3)

# 4 Objectives
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Research Objectives", "Evaluate practical efficiency under constrained and distributed hardware.")
objectives = [
    "Implement a single-GPU QLoRA pipeline for code instruction tuning under limited memory.",
    "Implement a dual-GPU DeepSpeed ZeRO-2 full fine-tuning pipeline for the same dataset and model family.",
    "Compare trainable parameters, runtime, throughput, loss, GPU memory, and utilization.",
    "Identify implementation challenges and explain when each approach is more suitable.",
]
for i, obj in enumerate(objectives):
    y = 1.45 + i * 1.25
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.95), Inches(y + 0.03), Inches(0.48), Inches(0.48))
    circ.fill.solid()
    circ.fill.fore_color.rgb = [TEAL, GOLD, BLUE, GREEN][i]
    circ.line.fill.background()
    n = slide.shapes.add_textbox(Inches(0.95), Inches(y + 0.11), Inches(0.48), Inches(0.25))
    npg = n.text_frame.paragraphs[0]
    npg.text = str(i + 1)
    npg.font.size = Pt(12)
    npg.font.bold = True
    npg.font.color.rgb = WHITE
    npg.alignment = PP_ALIGN.CENTER
    bullet_box(slide, 1.7, y, 10.4, 0.62, [obj], font_size=18, gap=False)
add_footer(slide, 4)

# 5 Proposed solution/workflow
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Proposed Solution", "A two-track workflow compares adapter-based tuning with distributed full fine-tuning.")
steps = [
    ("Load CodeAlpaca", "1,000 examples, seed 42"),
    ("Format Samples", "Question + reasoning + answer"),
    ("Tokenize", "Qwen tokenizer, max length 512"),
    ("Train Track A", "Single-GPU QLoRA"),
    ("Train Track B", "DeepSpeed ZeRO-2"),
    ("Compare", "Loss, memory, runtime, throughput"),
]
positions = [(0.75, 1.75), (2.85, 1.75), (4.95, 1.75), (2.05, 4.15), (5.55, 4.15), (9.0, 2.95)]
for idx, ((title_, body), (x, y)) in enumerate(zip(steps, positions)):
    add_card(slide, x, y, 1.85 if idx < 5 else 2.4, 1.05, title_, body, fill=WHITE, accent=[TEAL, BLUE, GREEN, TEAL, GOLD, RED][idx])
def connector(x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = MUTED
    c.line.width = Pt(1.4)
connector(2.6, 2.28, 2.85, 2.28)
connector(4.7, 2.28, 4.95, 2.28)
connector(5.88, 2.8, 2.98, 4.15)
connector(5.88, 2.8, 6.48, 4.15)
connector(3.9, 4.65, 9.0, 3.45)
connector(7.4, 4.65, 9.0, 3.75)
add_footer(slide, 5)

# 6 Methodology
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Methodology", "Both tracks used the same model and data preparation, then diverged by optimization strategy.")
rows = [
    ["Component", "QLoRA Track", "DeepSpeed ZeRO-2 Track"],
    ["Model", "Qwen/Qwen2.5-0.5B-Instruct", "Qwen/Qwen2.5-0.5B-Instruct"],
    ["Dataset", "CodeAlpaca 20K subset", "CodeAlpaca 20K subset"],
    ["Split", "900 train / 100 eval", "900 train / 100 eval"],
    ["Precision", "4-bit + FP16", "FP16"],
    ["Trainable scope", "LoRA adapters", "All parameters"],
    ["Hardware", "1 x Tesla T4", "2 x RTX A4000"],
    ["Epochs", "5", "1"],
]
add_table(slide, 0.75, 1.45, 11.85, 4.8, rows, col_widths=[2.2, 4.8, 4.8], font_size=10.5)
add_footer(slide, 6)

# 7 Challenges
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Implementation Challenges", "The project had practical systems constraints as much as modeling constraints.")
challenges = [
    ("Colab memory limit", "Full fine-tuning was not feasible on a single Tesla T4, forcing a QLoRA adaptation."),
    ("Distributed setup", "DeepSpeed required GPU verification, ZeRO config tuning, process launch, and metric collection."),
    ("Fair comparison", "The runs used different hardware and epoch counts, so throughput and loss must be interpreted carefully."),
    ("Telemetry mismatch", "nvidia-smi snapshots and PyTorch peak memory counters measure different memory behavior."),
    ("Data formatting", "CodeAlpaca has prompt-completion pairs, so a synthetic reasoning template was inserted consistently."),
]
for i, (t, b) in enumerate(challenges):
    x = 0.75 + (i % 2) * 6.0
    y = 1.35 + (i // 2) * 1.55
    add_card(slide, x, y, 5.45, 1.12, t, b, fill=LIGHT if i % 2 else WHITE, accent=[RED, GOLD, BLUE, TEAL, GREEN][i])
add_footer(slide, 7)

# 8 Baseline comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Baseline Comparison", "QLoRA serves as the low-resource baseline; ZeRO-2 is the distributed full-tuning approach.")
rows = [
    ["Metric", "QLoRA Baseline", "DeepSpeed ZeRO-2"],
    ["Trainable parameters", "8.80M", f"{deepspeed['trainable_params_m']:.2f}M"],
    ["Trainable percentage", "1.7497%", "100.0000%"],
    ["Runtime", "2,829.545 s", f"{deepspeed['runtime']:.3f} s"],
    ["Throughput", "1.590 samples/s", f"{deepspeed['samples_s']:.3f} samples/s"],
    ["Training loss", "0.448", f"{deepspeed['train_loss']:.3f}"],
    ["Evaluation loss", "0.823", f"{deepspeed['eval_loss']:.3f}"],
    ["Peak allocated memory", "1.66 GB", f"{deepspeed['peak_allocated_gb']:.4f} GB/rank"],
]
add_table(slide, 0.75, 1.35, 7.0, 5.35, rows, col_widths=[2.25, 2.35, 2.4], font_size=10.2)
add_picture(slide, params_chart, 8.12, 1.48, 4.42, 2.25)
add_picture(slide, memory_chart, 8.12, 4.15, 4.42, 2.25)
add_footer(slide, 8)

# 9 Results charts
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Graphical Results", "Runtime, throughput, and loss expose different efficiency profiles.")
add_picture(slide, runtime_chart, 0.65, 1.35, 3.95, 2.42)
add_picture(slide, throughput_chart, 4.72, 1.35, 3.95, 2.42)
add_picture(slide, loss_chart, 8.78, 1.35, 3.95, 2.42)
bullet_box(
    slide,
    0.9,
    4.45,
    11.4,
    1.45,
    [
        "DeepSpeed completed the training run about 11.45x faster in wall-clock runtime, but the QLoRA track ran for five epochs versus one.",
        "DeepSpeed produced higher sample throughput and lower evaluation loss; QLoRA produced lower training loss with far fewer trainable parameters.",
    ],
    font_size=16,
)
add_footer(slide, 9)

# 10 DeepSpeed telemetry
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "DeepSpeed Training Telemetry", "Logged loss and GPU summaries confirm the distributed run completed and used both GPUs.")
add_picture(slide, curve_chart, 0.78, 1.35, 5.65, 3.35)
add_picture(slide, gpu_chart, 6.72, 1.48, 5.65, 2.75)
add_card(slide, 6.8, 4.75, 2.55, 1.05, "Global steps", str(deepspeed["steps"]), fill=WHITE, accent=BLUE)
add_card(slide, 9.65, 4.75, 2.55, 1.05, "Eval samples/s", f"{evalm['eval_samples_per_second']:.3f}", fill=WHITE, accent=GREEN)
add_footer(slide, 10)

# 11 Interpretation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Result Interpretation", "The stronger approach depends on the resource constraint.")
add_card(slide, 0.82, 1.42, 5.4, 3.95, "When QLoRA Wins", "Best for low-cost experimentation, classroom environments, and cases where memory is the binding constraint. It updates only 1.7497% of parameters and fits comfortably on a Tesla T4.", fill=WHITE, accent=TEAL)
add_card(slide, 7.02, 1.42, 5.4, 3.95, "When DeepSpeed Wins", "Best when multiple GPUs are available and full-model adaptation is desired. It updates all parameters, improves throughput, and reached lower evaluation loss in the completed run.", fill=WHITE, accent=GOLD)
add_card(slide, 2.45, 5.72, 8.4, 0.7, "Key takeaway", "QLoRA optimizes accessibility; DeepSpeed ZeRO-2 optimizes distributed full fine-tuning performance.", fill=LIGHT, accent=BLUE)
add_footer(slide, 11)

# 12 Limitations and future work
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Limitations and Future Work", "The comparison is useful, but not hardware-neutral.")
bullet_box(
    slide,
    0.85,
    1.45,
    5.65,
    4.65,
    [
        "Different GPU environments: Tesla T4 on Colab versus two RTX A4000 GPUs on RunPod.",
        "Different epoch counts: QLoRA used five epochs; DeepSpeed used one epoch.",
        "Evaluation focused on loss and qualitative behavior, not HumanEval, MBPP, or a full code benchmark suite.",
        "Synthetic reasoning field may not represent verified reasoning traces.",
    ],
    font_size=16,
)
add_card(slide, 7.1, 1.55, 4.75, 1.2, "Future: Equalized Runs", "Repeat with matched epochs and a shared hardware budget.", fill=LIGHT, accent=BLUE)
add_card(slide, 7.1, 3.05, 4.75, 1.2, "Future: Larger Dataset", "Scale to 5,000 or 10,000 CodeAlpaca samples.", fill=LIGHT, accent=TEAL)
add_card(slide, 7.1, 4.55, 4.75, 1.2, "Future: Stronger Benchmarks", "Evaluate generated code with HumanEval or MBPP.", fill=LIGHT, accent=GOLD)
add_footer(slide, 12)

# 13 Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_title(slide, "Conclusion", "Efficient fine-tuning has two practical paths.")
add_metric(slide, 0.95, 1.65, "QLoRA trainable share", "1.75%", TEAL)
add_metric(slide, 4.05, 1.65, "DeepSpeed runtime", "247.1 s", GOLD)
add_metric(slide, 7.15, 1.65, "DeepSpeed throughput", "3.642/s", BLUE)
add_metric(slide, 10.25, 1.65, "DeepSpeed eval loss", "0.616", GREEN)
bullet_box(
    slide,
    1.0,
    3.6,
    11.0,
    1.7,
    [
        "QLoRA is the stronger choice when access is limited to a low-memory single GPU.",
        "DeepSpeed ZeRO-2 is better suited for faster distributed full-model fine-tuning when multi-GPU resources are available.",
        "The project demonstrates an end-to-end comparison using real training metrics, GPU telemetry, and reproducible experiment artifacts.",
    ],
    font_size=18,
)
add_footer(slide, 13)


out_path = OUT / "PDC_DeepSpeed_QLoRA_Presentation.pptx"
prs.save(out_path)
print(out_path)
