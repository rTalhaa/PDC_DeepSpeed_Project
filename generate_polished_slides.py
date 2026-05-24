from pathlib import Path
import json
import math
import random

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output" / "slides"
ASSETS = OUT / "assets_polished"
OUT.mkdir(parents=True, exist_ok=True)
ASSETS.mkdir(parents=True, exist_ok=True)

DEEPSPEED_RESULTS = ROOT / "DeepSpeed" / "results" / "metrics_zero2_fair.json"
GPU_SUMMARY = ROOT / "DeepSpeed" / "results" / "gpu_summary_zero2_fair.csv"
TRAINER_STATE = ROOT / "DeepSpeed" / "outputs" / "zero2_fair_run" / "trainer_state.json"

W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(12, 23, 38)
MIDNIGHT = RGBColor(15, 30, 50)
INK = RGBColor(28, 37, 49)
MUTED = RGBColor(100, 113, 128)
SUBTLE = RGBColor(229, 235, 241)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(31, 170, 155)
CYAN = RGBColor(61, 177, 222)
AMBER = RGBColor(239, 178, 68)
CORAL = RGBColor(220, 95, 89)
INDIGO = RGBColor(83, 112, 205)
GREEN = RGBColor(77, 164, 101)

HEX = {
    "navy": "#0c1726",
    "ink": "#1c2531",
    "muted": "#647180",
    "subtle": "#e5ebf1",
    "paper": "#f8fafc",
    "teal": "#1faa9b",
    "cyan": "#3db1de",
    "amber": "#efb244",
    "coral": "#dc5f59",
    "indigo": "#5370cd",
    "green": "#4da465",
}


qlora = {
    "hardware": "1 x Tesla T4",
    "method": "Single-GPU QLoRA",
    "trainable_params": 8_798_208,
    "total_params": 502_830_976,
    "trainable_pct": 1.7497,
    "epochs": 5,
    "steps": 565,
    "runtime": 2829.545,
    "samples_s": 1.590,
    "steps_s": 0.200,
    "train_loss": 0.448,
    "eval_loss": 0.823,
    "peak_allocated_gb": 1.66,
    "peak_reserved_gb": 3.56,
    "flops": 1.204e15,
}

with open(DEEPSPEED_RESULTS, "r", encoding="utf-8") as f:
    ds_metrics = json.load(f)
with open(TRAINER_STATE, "r", encoding="utf-8") as f:
    trainer_state = json.load(f)

train = ds_metrics["train_metrics"]
evalm = ds_metrics["eval_metrics"]
deepspeed = {
    "hardware": "2 x RTX A4000",
    "method": "Dual-GPU DeepSpeed ZeRO-2",
    "trainable_params": train["trainable_params"],
    "total_params": train["total_params"],
    "trainable_pct": train["trainable_percentage"],
    "epochs": int(train["epoch"]),
    "steps": trainer_state["global_step"],
    "runtime": train["train_runtime"],
    "samples_s": train["train_samples_per_second"],
    "steps_s": train["train_steps_per_second"],
    "train_loss": train["train_loss"],
    "eval_loss": evalm["eval_loss"],
    "peak_allocated_gb": train["rank_peak_memory_allocated_gb"],
    "peak_reserved_gb": train["rank_peak_memory_reserved_gb"],
    "flops": train["total_flos"],
    "eval_samples_s": evalm["eval_samples_per_second"],
}


def make_background(path, dark=True, seed=10):
    random.seed(seed)
    img = Image.new("RGB", (1920, 1080), "#0c1726" if dark else "#f8fafc")
    px = img.load()
    for y in range(img.height):
        for x in range(img.width):
            if dark:
                r = 10 + int(12 * x / img.width)
                g = 22 + int(22 * y / img.height)
                b = 38 + int(36 * (x + y) / (img.width + img.height))
            else:
                r = 248 - int(7 * y / img.height)
                g = 250 - int(8 * x / img.width)
                b = 252 - int(8 * (x + y) / (img.width + img.height))
            px[x, y] = (r, g, b)
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    if dark:
        for i in range(34):
            x = random.randint(0, 1920)
            y = random.randint(0, 1080)
            radius = random.randint(2, 5)
            color = random.choice([(31, 170, 155, 180), (61, 177, 222, 165), (239, 178, 68, 145)])
            d.ellipse((x - radius, y - radius, x + radius, y + radius), fill=color)
        for i in range(18):
            x1 = random.randint(0, 1920)
            y1 = random.randint(0, 1080)
            x2 = x1 + random.randint(-360, 360)
            y2 = y1 + random.randint(-240, 240)
            d.line((x1, y1, x2, y2), fill=(61, 177, 222, 60), width=2)
        d.polygon([(1180, -40), (1940, 0), (1940, 1080), (1510, 1080)], fill=(31, 170, 155, 34))
        d.polygon([(1390, -40), (1940, 0), (1940, 830), (1710, 1080)], fill=(239, 178, 68, 32))
        for x in range(1040, 1840, 120):
            d.line((x, 50, x + 430, 1010), fill=(255, 255, 255, 18), width=1)
        for y in range(80, 1060, 110):
            d.line((1060, y, 1910, y + 150), fill=(255, 255, 255, 16), width=1)
    else:
        d.ellipse((1340, -260, 2120, 520), fill=(31, 170, 155, 30))
        d.ellipse((-260, 560, 460, 1280), fill=(83, 112, 205, 22))
        d.rectangle((0, 0, 1920, 14), fill=(31, 170, 155, 255))
    img = Image.alpha_composite(img.convert("RGBA"), overlay).filter(ImageFilter.SMOOTH)
    img.save(path)
    return path


cover_bg = make_background(ASSETS / "cover_bg.png", dark=True, seed=19)
light_bg = make_background(ASSETS / "light_bg.png", dark=False, seed=21)


def chart_style(ax):
    ax.set_facecolor("white")
    ax.grid(axis="y", color=HEX["subtle"], linewidth=0.9)
    ax.set_axisbelow(True)
    ax.tick_params(colors=HEX["muted"], labelsize=9)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color(HEX["subtle"])
    ax.spines["bottom"].set_color(HEX["subtle"])


def save_bar_chart(filename, title, labels, values, ylabel, colors, fmt="{:.2f}", horizontal=False):
    fig, ax = plt.subplots(figsize=(7.4, 4.2), dpi=190)
    if horizontal:
        bars = ax.barh(labels, values, color=colors, height=0.48)
        ax.invert_yaxis()
        ax.set_xlabel(ylabel, color=HEX["muted"])
        xmax = max(values) * 1.22
        ax.set_xlim(0, xmax)
        for b, v in zip(bars, values):
            ax.text(v + xmax * 0.025, b.get_y() + b.get_height() / 2, fmt.format(v), va="center", fontsize=10, color=HEX["ink"], fontweight="bold")
    else:
        bars = ax.bar(labels, values, color=colors, width=0.52)
        ax.set_ylabel(ylabel, color=HEX["muted"])
        ymax = max(values) * 1.22
        ax.set_ylim(0, ymax)
        for b, v in zip(bars, values):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + ymax * 0.025, fmt.format(v), ha="center", va="bottom", fontsize=10, color=HEX["ink"], fontweight="bold")
    chart_style(ax)
    ax.set_title(title, loc="left", fontsize=15, fontweight="bold", color=HEX["ink"], pad=12)
    fig.tight_layout(pad=1.4)
    path = ASSETS / filename
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def save_loss_chart():
    fig, ax = plt.subplots(figsize=(7.4, 4.2), dpi=190)
    x = np.arange(2)
    width = 0.32
    ax.bar(x - width / 2, [qlora["train_loss"], qlora["eval_loss"]], width, label="QLoRA", color=HEX["teal"])
    ax.bar(x + width / 2, [deepspeed["train_loss"], deepspeed["eval_loss"]], width, label="DeepSpeed", color=HEX["amber"])
    ax.set_xticks(x, ["Training", "Evaluation"])
    ax.set_ylabel("Loss", color=HEX["muted"])
    ax.set_ylim(0, 1.03)
    chart_style(ax)
    ax.set_title("Loss Comparison", loc="left", fontsize=15, fontweight="bold", color=HEX["ink"], pad=12)
    ax.legend(frameon=False, loc="upper right")
    for c in ax.containers:
        ax.bar_label(c, fmt="%.3f", padding=3, fontsize=9, color=HEX["ink"])
    fig.tight_layout(pad=1.4)
    path = ASSETS / "loss_comparison.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def save_curve_chart():
    history = [h for h in trainer_state["log_history"] if "loss" in h]
    steps = [h["step"] for h in history]
    losses = [h["loss"] for h in history]
    fig, ax = plt.subplots(figsize=(7.4, 4.2), dpi=190)
    ax.plot(steps, losses, color=HEX["indigo"], marker="o", markersize=5, linewidth=2.5)
    ax.fill_between(steps, losses, [min(losses) * 0.92] * len(losses), color=HEX["indigo"], alpha=0.10)
    ax.axhline(deepspeed["eval_loss"], linestyle="--", color=HEX["coral"], linewidth=1.8, label=f"Eval loss {deepspeed['eval_loss']:.3f}")
    ax.set_xlabel("Global step", color=HEX["muted"])
    ax.set_ylabel("Loss", color=HEX["muted"])
    chart_style(ax)
    ax.set_title("DeepSpeed ZeRO-2 Logged Training Loss", loc="left", fontsize=15, fontweight="bold", color=HEX["ink"], pad=12)
    ax.legend(frameon=False)
    fig.tight_layout(pad=1.4)
    path = ASSETS / "deepspeed_curve.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def save_gpu_chart():
    df = pd.read_csv(GPU_SUMMARY)
    labels = [f"GPU {int(v)}" for v in df["index"]]
    fig, ax = plt.subplots(figsize=(7.4, 4.2), dpi=190)
    x = np.arange(len(labels))
    width = 0.34
    ax.bar(x - width / 2, df["avg_gpu_util_percent"], width, label="Avg utilization (%)", color=HEX["cyan"])
    ax.bar(x + width / 2, df["max_gpu_util_percent"], width, label="Max utilization (%)", color=HEX["teal"])
    ax.set_xticks(x, labels)
    ax.set_ylim(0, 115)
    chart_style(ax)
    ax.set_title("Dual-GPU Utilization Summary", loc="left", fontsize=15, fontweight="bold", color=HEX["ink"], pad=12)
    ax.legend(frameon=False)
    for c in ax.containers:
        ax.bar_label(c, fmt="%.0f", padding=3, fontsize=9)
    fig.tight_layout(pad=1.4)
    path = ASSETS / "gpu_utilization.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def save_tradeoff_bubble():
    fig, ax = plt.subplots(figsize=(7.4, 4.2), dpi=190)
    xs = [qlora["samples_s"], deepspeed["samples_s"]]
    ys = [qlora["eval_loss"], deepspeed["eval_loss"]]
    sizes = [qlora["peak_allocated_gb"] * 520, deepspeed["peak_allocated_gb"] * 520]
    ax.scatter(xs, ys, s=sizes, color=[HEX["teal"], HEX["amber"]], alpha=0.78, edgecolor="white", linewidth=2)
    ax.text(xs[0] + 0.07, ys[0] + 0.025, "QLoRA", fontsize=11, fontweight="bold", color=HEX["ink"])
    ax.text(xs[1] - 0.62, ys[1] - 0.055, "DeepSpeed", fontsize=11, fontweight="bold", color=HEX["ink"])
    ax.set_xlabel("Training samples / second", color=HEX["muted"])
    ax.set_ylabel("Evaluation loss", color=HEX["muted"])
    ax.set_xlim(1.0, 4.2)
    ax.set_ylim(0.55, 0.9)
    chart_style(ax)
    ax.set_title("Performance Trade-off", loc="left", fontsize=15, fontweight="bold", color=HEX["ink"], pad=12)
    ax.text(1.02, 0.565, "Bubble size = peak allocated GPU memory", fontsize=8.5, color=HEX["muted"])
    fig.tight_layout(pad=1.4)
    path = ASSETS / "tradeoff_bubble.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


runtime_chart = save_bar_chart("runtime.png", "Training Runtime", ["QLoRA", "DeepSpeed"], [qlora["runtime"], deepspeed["runtime"]], "Seconds", [HEX["teal"], HEX["amber"]], fmt="{:,.0f}")
throughput_chart = save_bar_chart("throughput.png", "Training Throughput", ["QLoRA", "DeepSpeed"], [qlora["samples_s"], deepspeed["samples_s"]], "Samples / second", [HEX["teal"], HEX["amber"]])
params_chart = save_bar_chart("params.png", "Trainable Parameter Share", ["QLoRA", "DeepSpeed"], [qlora["trainable_pct"], deepspeed["trainable_pct"]], "Percent", [HEX["teal"], HEX["amber"]], fmt="{:.1f}%")
memory_chart = save_bar_chart("memory.png", "Peak Allocated Memory", ["QLoRA", "DeepSpeed"], [qlora["peak_allocated_gb"], deepspeed["peak_allocated_gb"]], "GB", [HEX["teal"], HEX["amber"]])
loss_chart = save_loss_chart()
curve_chart = save_curve_chart()
gpu_chart = save_gpu_chart()
tradeoff_chart = save_tradeoff_bubble()


prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def rgb_tuple(c):
    return (c[0], c[1], c[2])


def bg(slide, fill=PAPER, image=None):
    if image:
        slide.shapes.add_picture(str(image), 0, 0, W, H)
    else:
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
        r.fill.solid()
        r.fill.fore_color.rgb = fill
        r.line.fill.background()
        slide.shapes._spTree.remove(r._element)
        slide.shapes._spTree.insert(2, r._element)


def textbox(slide, x, y, w, h, text, size=16, color=INK, bold=False, font="Aptos", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def title(slide, text, kicker=None, dark=False):
    color = WHITE if dark else NAVY
    muted = RGBColor(190, 204, 218) if dark else MUTED
    if kicker:
        textbox(slide, 0.72, 0.36, 11.8, 0.28, kicker.upper(), 8.5, muted, bold=True)
    textbox(slide, 0.72, 0.62, 10.9, 0.65, text, 27, color, bold=True, font="Aptos Display")
    if not dark:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.25), Inches(1.18), Inches(0.045))
        line.fill.solid()
        line.fill.fore_color.rgb = TEAL
        line.line.fill.background()


def footer(slide, n, dark=False):
    c = RGBColor(173, 188, 203) if dark else MUTED
    linec = RGBColor(62, 80, 101) if dark else SUBTLE
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(7.05), Inches(11.85), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = linec
    line.line.fill.background()
    textbox(slide, 0.72, 7.12, 6.0, 0.2, "PDC DeepSpeed Project | QLoRA vs ZeRO-2", 8.5, c)
    textbox(slide, 12.0, 7.12, 0.55, 0.2, f"{n:02d}", 8.5, c, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=WHITE, line=SUBTLE, shadow=False):
    if shadow:
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.035), Inches(y + 0.045), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(218, 226, 235)
        s.line.fill.background()
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    return shp


def accent_card(slide, x, y, w, h, head, body, accent=TEAL):
    card(slide, x, y, w, h, WHITE, SUBTLE, shadow=True)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.09), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x + 0.28, y + 0.22, w - 0.45, 0.32, head, 14, NAVY, bold=True, font="Aptos Display")
    textbox(slide, x + 0.28, y + 0.68, w - 0.45, h - 0.78, body, 11.5, MUTED)


def bullet_list(slide, x, y, w, h, items, size=15, color=INK, dot=TEAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
        p.level = 0
    return box


def metric_tile(slide, x, y, w, h, label, value, sub="", accent=TEAL, dark=False):
    fill = RGBColor(20, 38, 61) if dark else WHITE
    line = RGBColor(52, 75, 99) if dark else SUBTLE
    card(slide, x, y, w, h, fill, line, shadow=not dark)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.06))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    label_color = RGBColor(179, 195, 210) if dark else MUTED
    value_color = WHITE if dark else NAVY
    textbox(slide, x + 0.20, y + 0.20, w - 0.4, 0.22, label.upper(), 8, label_color, bold=True)
    textbox(slide, x + 0.20, y + 0.48, w - 0.4, 0.48, value, 23, value_color, bold=True, font="Aptos Display")
    if sub:
        textbox(slide, x + 0.20, y + 0.98, w - 0.4, 0.25, sub, 9, label_color)


def add_table(slide, x, y, w, h, rows, widths, size=9.5):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(text)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (WHITE if r % 2 else PAPER)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(size)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else INK
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    return shape


def pic(slide, path, x, y, w, h):
    card(slide, x - 0.02, y - 0.02, w + 0.04, h + 0.04, WHITE, SUBTLE, shadow=True)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def chip_icon(slide, x, y, size, color):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    outer.fill.solid()
    outer.fill.fore_color.rgb = color
    outer.line.fill.background()
    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + size * 0.27), Inches(y + size * 0.27), Inches(size * 0.46), Inches(size * 0.46))
    inner.fill.solid()
    inner.fill.fore_color.rgb = WHITE
    inner.fill.transparency = 15
    inner.line.fill.background()
    for i in range(4):
        offset = size * (0.18 + i * 0.18)
        for side in [0, 1]:
            pin = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x - 0.05 if side == 0 else x + size), Inches(y + offset), Inches(0.05), Inches(0.025))
            pin.fill.solid()
            pin.fill.fore_color.rgb = color
            pin.line.fill.background()


def connector(slide, x1, y1, x2, y2, color=MUTED):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(1.5)
    return c


# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=cover_bg)
textbox(s, 0.78, 0.68, 1.85, 0.32, "FINAL PRESENTATION", 9, RGBColor(188, 205, 220), bold=True)
textbox(s, 0.78, 1.12, 8.7, 1.55, "Efficient Fine-Tuning of Code Language Models", 37, WHITE, bold=True, font="Aptos Display")
textbox(s, 0.82, 2.86, 8.65, 0.55, "Single-GPU QLoRA vs Dual-GPU DeepSpeed ZeRO-2 under hardware constraints", 18, RGBColor(202, 216, 229))
metric_tile(s, 0.82, 4.55, 2.45, 1.25, "Model", "Qwen2.5", "0.5B-Instruct", TEAL, dark=True)
metric_tile(s, 3.55, 4.55, 2.45, 1.25, "Dataset", "CodeAlpaca", "1,000 sample subset", AMBER, dark=True)
metric_tile(s, 6.28, 4.55, 2.45, 1.25, "Split", "900 / 100", "train / evaluation", CYAN, dark=True)
metric_tile(s, 9.01, 4.55, 2.45, 1.25, "Question", "Efficiency", "memory, speed, loss", GREEN, dark=True)
textbox(s, 0.82, 6.45, 5.8, 0.32, "Talha Rashid, Bilal Fayyaz | Department of Computer Science", 10.5, RGBColor(190, 204, 218))
footer(s, 1, dark=True)

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "The Problem", "basic introduction")
textbox(s, 0.82, 1.58, 5.6, 0.65, "Fine-tuning code LLMs is valuable, but the hardware bill arrives before the model improves.", 21, NAVY, bold=True, font="Aptos Display")
bullet_list(s, 0.92, 2.58, 5.35, 2.8, [
    "Full-parameter training stresses GPU memory and runtime.",
    "Single cloud notebook GPUs often cannot support full fine-tuning comfortably.",
    "Distributed training improves throughput, but adds setup and synchronization complexity.",
    "The project compares two practical efficiency paths on the same task."
], 15)
for i, (head, val, c) in enumerate([("Memory", "Can the model fit?", TEAL), ("Speed", "How fast does it train?", AMBER), ("Quality", "What loss is achieved?", INDIGO)]):
    accent_card(s, 7.0, 1.55 + i * 1.45, 4.95, 1.0, head, val, c)
footer(s, 2)

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Literature Review", "base paper highlighted")
accent_card(s, 0.82, 1.45, 3.75, 2.2, "Base Paper", "The provided IEEE paper compares Single-GPU QLoRA and Dual-GPU DeepSpeed ZeRO-2 for Qwen2.5-0.5B-Instruct on CodeAlpaca.", INDIGO)
accent_card(s, 4.82, 1.45, 3.55, 2.2, "LoRA / QLoRA", "LoRA trains low-rank adapters. QLoRA loads the base model in 4-bit form and backpropagates into adapters.", TEAL)
accent_card(s, 8.62, 1.45, 3.55, 2.2, "DeepSpeed ZeRO", "ZeRO partitions optimizer states and gradients across workers; ZeRO-2 enables practical distributed full fine-tuning.", AMBER)
accent_card(s, 0.82, 4.15, 3.75, 1.45, "CodeAlpaca", "Prompt-completion code instruction examples are reformatted into question, reasoning, and answer blocks.", GREEN)
accent_card(s, 4.82, 4.15, 3.55, 1.45, "Hugging Face Stack", "Transformers, Trainer, PEFT-style workflows, tokenizer handling, and metric capture support reproducibility.", CYAN)
accent_card(s, 8.62, 4.15, 3.55, 1.45, "Research Gap", "A practical comparison helps choose between accessible adapter tuning and faster distributed full-model updates.", CORAL)
footer(s, 3)

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Research Objectives", "what the study measures")
objs = [
    ("Implement QLoRA", "Fine-tune Qwen2.5-0.5B-Instruct on a single Tesla T4 using 4-bit quantization and LoRA adapters.", TEAL),
    ("Implement ZeRO-2", "Run full-parameter fine-tuning with DeepSpeed across two RTX A4000 GPUs.", AMBER),
    ("Measure Trade-offs", "Compare trainable parameters, runtime, throughput, loss, memory, FLOPs, and GPU utilization.", INDIGO),
    ("Recommend Usage", "Explain which method is better under constrained single-GPU versus multi-GPU conditions.", GREEN),
]
for i, (h, b, c) in enumerate(objs):
    x = 0.82 + (i % 2) * 6.0
    y = 1.55 + (i // 2) * 2.05
    card(s, x, y, 5.3, 1.42, WHITE, SUBTLE, shadow=True)
    chip_icon(s, x + 0.25, y + 0.32, 0.55, c)
    textbox(s, x + 1.02, y + 0.22, 4.05, 0.32, h, 16, NAVY, bold=True, font="Aptos Display")
    textbox(s, x + 1.02, y + 0.68, 4.05, 0.5, b, 11.5, MUTED)
footer(s, 4)

# Slide 5 workflow
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Proposed Solution Workflow", "methodology diagram")
textbox(s, 0.88, 1.45, 5.7, 0.35, "Shared preparation", 13, NAVY, bold=True)
shared = [("Load", "CodeAlpaca 20K"), ("Sample", "1,000 examples"), ("Format", "Q + R + A"), ("Tokenize", "max length 512")]
for i, (h, b) in enumerate(shared):
    x = 0.88 + i * 2.12
    accent_card(s, x, 1.95, 1.66, 0.92, h, b, [TEAL, CYAN, GREEN, INDIGO][i])
    if i < len(shared) - 1:
        connector(s, x + 1.66, 2.42, x + 2.08, 2.42, RGBColor(138, 153, 169))
connector(s, 4.35, 2.88, 3.15, 3.78, RGBColor(138, 153, 169))
connector(s, 4.35, 2.88, 7.28, 3.78, RGBColor(138, 153, 169))
accent_card(s, 1.18, 3.78, 4.2, 1.42, "Track A: QLoRA Baseline", "4-bit loading, LoRA adapters, FP16, gradient accumulation, single Tesla T4.", TEAL)
accent_card(s, 6.05, 3.78, 4.2, 1.42, "Track B: DeepSpeed ZeRO-2", "FP16 full fine-tuning, optimizer and gradient partitioning, two RTX A4000 GPUs.", AMBER)
connector(s, 5.38, 4.48, 10.85, 4.48, RGBColor(138, 153, 169))
connector(s, 10.25, 4.48, 10.85, 4.48, RGBColor(138, 153, 169))
accent_card(s, 10.65, 3.5, 1.85, 1.95, "Compare", "loss, runtime, throughput, memory, utilization", CORAL)
footer(s, 5)

# Slide 6 methodology table
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Methodology", "same task, different optimization strategy")
rows = [
    ["Component", "Single-GPU QLoRA", "Dual-GPU DeepSpeed ZeRO-2"],
    ["Platform", "Google Colab", "RunPod"],
    ["GPU", "1 x Tesla T4", "2 x RTX A4000"],
    ["Training type", "4-bit LoRA adapter tuning", "Full-parameter fine-tuning"],
    ["Precision", "4-bit + FP16", "FP16"],
    ["Dataset split", "900 train / 100 eval", "900 train / 100 eval"],
    ["Max sequence length", "512", "512"],
    ["Epochs", "5", "1"],
    ["Global steps", "565", "57"],
]
add_table(s, 0.82, 1.48, 11.75, 4.95, rows, [2.5, 4.55, 4.7], 10)
textbox(s, 0.85, 6.62, 10.8, 0.25, "Source: paper Table I and methodology section.", 8.8, MUTED)
footer(s, 6)

# Slide 7 fact checked setup
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Fact-Checked Experimental Setup", "aligned to the provided paper")
metric_tile(s, 0.9, 1.55, 2.65, 1.25, "Dataset", "CodeAlpaca 20K", "HuggingFaceH4 subset", TEAL)
metric_tile(s, 3.85, 1.55, 2.65, 1.25, "Samples", "1,000", "900 train, 100 eval", CYAN)
metric_tile(s, 6.8, 1.55, 2.65, 1.25, "Model", "Qwen2.5-0.5B", "Instruct checkpoint", AMBER)
metric_tile(s, 9.75, 1.55, 2.65, 1.25, "Sequence", "512", "maximum tokens", INDIGO)
accent_card(s, 1.05, 3.55, 4.95, 1.4, "QLoRA configuration", "LoRA rank 16, alpha 32, dropout 0.05; target modules include q_proj, k_proj, v_proj, o_proj, gate_proj, up_proj, and down_proj.", TEAL)
accent_card(s, 7.08, 3.55, 4.95, 1.4, "ZeRO-2 configuration", "FP16 training with gradient accumulation, reduce-scatter, all-gather partitioning, and contiguous gradients.", AMBER)
textbox(s, 1.05, 5.72, 10.8, 0.5, "The earlier deck's core values were checked against the paper. This version keeps the paper's final reported QLoRA runtime of 2,829.545 s and losses of 0.448 / 0.823.", 12, MUTED)
footer(s, 7)

# Slide 8 challenges
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Implementation Challenges", "what made the project non-trivial")
challenges = [
    ("Single-GPU limitation", "Colab's Tesla T4 could not comfortably support full fine-tuning, so QLoRA was used for the constrained track.", TEAL),
    ("Distributed launch complexity", "DeepSpeed required multi-GPU verification, config tuning, launcher commands, and synchronized metric collection.", AMBER),
    ("Comparison fairness", "Hardware and epoch counts differ, so runtime and loss are interpreted as practical outcomes rather than a controlled benchmark.", INDIGO),
    ("Telemetry interpretation", "nvidia-smi snapshots and PyTorch peak memory report different views of memory behavior.", CYAN),
    ("Dataset structure", "CodeAlpaca provides prompt-completion pairs, so a consistent reasoning template was inserted for the desired response format.", GREEN),
]
for i, (h, b, c) in enumerate(challenges):
    x = 0.82 + (i % 2) * 6.0
    y = 1.42 + (i // 2) * 1.62
    accent_card(s, x, y, 5.35, 1.08, h, b, c)
footer(s, 8)

# Slide 9 comparison table
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Baseline Comparison", "QLoRA baseline versus proposed distributed full fine-tuning")
rows = [
    ["Metric", "QLoRA Baseline", "DeepSpeed ZeRO-2"],
    ["Trainable parameters", "8,798,208", f"{deepspeed['trainable_params']:,}"],
    ["Trainable percentage", "1.7497%", "100.0000%"],
    ["Training runtime", "2,829.545 s", f"{deepspeed['runtime']:.3f} s"],
    ["Training throughput", "1.590 samples/s", f"{deepspeed['samples_s']:.3f} samples/s"],
    ["Training loss", "0.448", f"{deepspeed['train_loss']:.3f}"],
    ["Evaluation loss", "0.823", f"{deepspeed['eval_loss']:.3f}"],
    ["Peak allocated memory", "1.66 GB", f"{deepspeed['peak_allocated_gb']:.4f} GB per rank"],
]
add_table(s, 0.82, 1.42, 7.25, 4.95, rows, [2.55, 2.35, 2.35], 9.7)
pic(s, params_chart, 8.45, 1.5, 3.85, 2.25)
pic(s, memory_chart, 8.45, 4.18, 3.85, 2.05)
footer(s, 9)

# Slide 10 graphical results
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Graphical Results", "speed and loss at a glance")
pic(s, runtime_chart, 0.76, 1.42, 3.75, 2.35)
pic(s, throughput_chart, 4.80, 1.42, 3.75, 2.35)
pic(s, loss_chart, 8.84, 1.42, 3.75, 2.35)
accent_card(s, 0.95, 4.55, 3.6, 1.2, "Runtime", "DeepSpeed finished the one-epoch run in 247.091 s versus 2,829.545 s for the five-epoch QLoRA run.", AMBER)
accent_card(s, 4.88, 4.55, 3.6, 1.2, "Throughput", "DeepSpeed reached 3.642 samples/s compared with 1.590 samples/s for QLoRA.", CYAN)
accent_card(s, 8.80, 4.55, 3.6, 1.2, "Loss", "QLoRA had lower training loss; DeepSpeed had lower evaluation loss after full-parameter tuning.", GREEN)
footer(s, 10)

# Slide 11 telemetry
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "DeepSpeed Telemetry", "loss trend and GPU utilization")
pic(s, curve_chart, 0.8, 1.45, 5.65, 3.35)
pic(s, gpu_chart, 6.82, 1.45, 5.35, 3.35)
metric_tile(s, 1.05, 5.45, 2.55, 1.0, "Global steps", "57", "one epoch", INDIGO)
metric_tile(s, 3.95, 5.45, 2.55, 1.0, "Eval speed", f"{deepspeed['eval_samples_s']:.1f}/s", "samples per second", GREEN)
metric_tile(s, 6.85, 5.45, 2.55, 1.0, "Max util", "100%", "GPU 1 peak", TEAL)
metric_tile(s, 9.75, 5.45, 2.55, 1.0, "Avg power", "~67 W", "per GPU", AMBER)
footer(s, 11)

# Slide 12 tradeoff
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Performance Trade-off", "the better approach depends on the constraint")
pic(s, tradeoff_chart, 0.88, 1.48, 5.85, 3.45)
accent_card(s, 7.25, 1.55, 4.65, 1.3, "QLoRA optimizes accessibility", "It updates only 1.7497% of parameters and stays within a low memory footprint on a single Tesla T4.", TEAL)
accent_card(s, 7.25, 3.15, 4.65, 1.3, "DeepSpeed optimizes full adaptation", "It updates 100.0000% of parameters and reaches stronger throughput and lower evaluation loss with two GPUs.", AMBER)
accent_card(s, 7.25, 4.75, 4.65, 1.3, "Interpret carefully", "The runs are practical comparisons, not perfectly controlled hardware-neutral benchmarks.", CORAL)
footer(s, 12)

# Slide 13 limitations
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=light_bg)
title(s, "Limitations and Future Work", "how to make the comparison stronger")
bullet_list(s, 0.95, 1.5, 5.55, 4.25, [
    "Different hardware: Tesla T4 versus two RTX A4000 GPUs.",
    "Different epoch counts: QLoRA used five epochs, DeepSpeed used one.",
    "Reasoning text was synthetically inserted because CodeAlpaca contains prompt-completion pairs.",
    "Evaluation used loss and limited qualitative testing, not HumanEval or MBPP."
], 15)
accent_card(s, 7.15, 1.55, 4.65, 1.15, "Next experiment", "Run matched epoch counts on equalized hardware budgets.", INDIGO)
accent_card(s, 7.15, 3.05, 4.65, 1.15, "Scale data", "Repeat with 5,000 or 10,000 CodeAlpaca examples.", TEAL)
accent_card(s, 7.15, 4.55, 4.65, 1.15, "Add benchmarks", "Evaluate generated code using HumanEval, MBPP, or similar suites.", AMBER)
footer(s, 13)

# Slide 14 conclusion
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, image=cover_bg)
title(s, "Conclusion", "final takeaway", dark=True)
metric_tile(s, 0.85, 1.65, 2.55, 1.2, "QLoRA", "1.75%", "trainable parameters", TEAL, dark=True)
metric_tile(s, 3.7, 1.65, 2.55, 1.2, "DeepSpeed", "247.1 s", "training runtime", AMBER, dark=True)
metric_tile(s, 6.55, 1.65, 2.55, 1.2, "Throughput", "3.642/s", "DeepSpeed samples", CYAN, dark=True)
metric_tile(s, 9.4, 1.65, 2.55, 1.2, "Eval loss", "0.616", "DeepSpeed result", GREEN, dark=True)
textbox(s, 0.92, 3.62, 10.5, 0.58, "QLoRA is the stronger choice for constrained single-GPU environments; DeepSpeed ZeRO-2 is better when multi-GPU resources make full-model adaptation practical.", 21, WHITE, bold=True, font="Aptos Display")
textbox(s, 0.95, 4.72, 10.9, 0.72, "The project demonstrates a reproducible, metric-backed comparison of memory efficiency, throughput, loss, and implementation complexity.", 15, RGBColor(204, 218, 231))
footer(s, 14, dark=True)


out_path = OUT / "PDC_DeepSpeed_QLoRA_Presentation_POLISHED.pptx"
prs.save(out_path)
print(out_path)
